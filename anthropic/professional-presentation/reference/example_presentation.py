#!/usr/bin/env python3
"""
Example: Creating a Professional Business Presentation with Pastel Colors

This example demonstrates how to create a complete PowerPoint presentation
with a modern pastel color scheme. Perfect for any corporate presentation.

Usage:
    python3 example_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import os

# Professional Pastel Color Palette
PASTEL_BLUE = RGBColor(147, 197, 253)
PASTEL_LAVENDER = RGBColor(196, 181, 253)
PASTEL_MINT = RGBColor(167, 243, 208)
PASTEL_PEACH = RGBColor(254, 202, 202)
PASTEL_CORAL = RGBColor(252, 165, 165)
PASTEL_SAGE = RGBColor(187, 222, 214)
PASTEL_PERIWINKLE = RGBColor(199, 210, 254)
PASTEL_ROSE = RGBColor(251, 207, 232)

# Text colors
DARK_SLATE = RGBColor(51, 65, 85)
MEDIUM_SLATE = RGBColor(100, 116, 139)
LIGHT_GRAY = RGBColor(241, 245, 249)


def create_title_slide(prs, title, subtitle=""):
    """Create professional title slide with pastel colors"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = PASTEL_BLUE
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add subtitle if provided
    if subtitle:
        left = Inches(1)
        top = Inches(3.2)
        width = Inches(8)
        height = Inches(0.8)

        subtitle_box = slide.shapes.add_textbox(left, top, width, height)
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].font.size = Pt(24)
        subtitle_frame.paragraphs[0].font.color.rgb = MEDIUM_SLATE
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add decorative pastel rectangle at bottom
    left = Inches(3)
    top = Inches(4.8)
    width = Inches(4)
    height = Inches(0.1)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PASTEL_LAVENDER
    shape.line.color.rgb = PASTEL_LAVENDER

    return slide


def create_section_header(prs, section_title):
    """Create section divider slide with pastel background"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add light pastel background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_GRAY

    # Add centered title
    left = Inches(1)
    top = Inches(2.3)
    width = Inches(8)
    height = Inches(1)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = section_title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = PASTEL_BLUE
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide


def create_content_slide(prs, title, bullet_points):
    """
    Create content slide with bullet points
    bullet_points: list of tuples [(level, text), ...]
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = DARK_SLATE

    # Add pastel accent line under title
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(2)
    height = Inches(0.05)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PASTEL_BLUE
    shape.line.color.rgb = PASTEL_BLUE

    # Add content area for bullets
    left = Inches(0.8)
    top = Inches(1.6)
    width = Inches(8.5)
    height = Inches(3.5)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    for idx, (level, text) in enumerate(bullet_points):
        if idx == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = text
        p.level = level
        p.font.size = Pt(18 - (level * 2))
        p.font.color.rgb = DARK_SLATE
        p.space_before = Pt(6) if level == 0 else Pt(3)

    return slide


def create_two_column_slide(prs, title, left_title, left_points, right_title, right_points):
    """Create two-column comparison slide with pastel accents"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = DARK_SLATE

    # Left column with pastel blue accent
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(4.3)
    height = Inches(3.8)

    left_box = slide.shapes.add_shape(1, left, top, width, height)
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
    left_box.line.color.rgb = PASTEL_BLUE
    left_box.line.width = Pt(2)

    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(4), Inches(0.5))
    left_title_frame = left_title_box.text_frame
    left_title_frame.text = left_title
    left_title_frame.paragraphs[0].font.size = Pt(20)
    left_title_frame.paragraphs[0].font.bold = True
    left_title_frame.paragraphs[0].font.color.rgb = PASTEL_BLUE

    # Left column content
    left_content_box = slide.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(3.6), Inches(2.8))
    left_content_frame = left_content_box.text_frame
    for idx, point in enumerate(left_points):
        if idx == 0:
            p = left_content_frame.paragraphs[0]
        else:
            p = left_content_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_SLATE
        p.space_before = Pt(6)

    # Right column with pastel mint accent
    left = Inches(5.2)
    top = Inches(1.5)
    width = Inches(4.3)
    height = Inches(3.8)

    right_box = slide.shapes.add_shape(1, left, top, width, height)
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(240, 253, 244)
    right_box.line.color.rgb = PASTEL_MINT
    right_box.line.width = Pt(2)

    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(5.4), Inches(1.7), Inches(4), Inches(0.5))
    right_title_frame = right_title_box.text_frame
    right_title_frame.text = right_title
    right_title_frame.paragraphs[0].font.size = Pt(20)
    right_title_frame.paragraphs[0].font.bold = True
    right_title_frame.paragraphs[0].font.color.rgb = PASTEL_MINT

    # Right column content
    right_content_box = slide.shapes.add_textbox(Inches(5.6), Inches(2.3), Inches(3.6), Inches(2.8))
    right_content_frame = right_content_box.text_frame
    for idx, point in enumerate(right_points):
        if idx == 0:
            p = right_content_frame.paragraphs[0]
        else:
            p = right_content_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_SLATE
        p.space_before = Pt(6)

    return slide


def create_table_slide(prs, title, table_data):
    """Create table slide with pastel header"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = DARK_SLATE

    # Create table
    rows = len(table_data)
    cols = len(table_data[0])

    left = Inches(0.8)
    top = Inches(1.6)
    width = Inches(8.5)
    height = Inches(3.5)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Format table
    for i, row_data in enumerate(table_data):
        for j, cell_value in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_value)

            # Format header row
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PASTEL_BLUE
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.bold = True
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = DARK_SLATE
            else:
                # Alternating row colors
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_GRAY
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = DARK_SLATE

            paragraph.alignment = PP_ALIGN.CENTER

    return slide


def create_bar_chart_slide(prs, title, categories, series_data):
    """Create bar chart with pastel colors"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = DARK_SLATE

    # Add chart
    chart_data = CategoryChartData()
    chart_data.categories = categories

    for series_name, values in series_data:
        chart_data.add_series(series_name, values)

    x, y, cx, cy = Inches(1), Inches(1.6), Inches(8), Inches(3.8)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    # Style chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM

    # Apply pastel colors to series
    pastel_colors = [
        PASTEL_BLUE,
        PASTEL_MINT,
        PASTEL_LAVENDER,
        PASTEL_PEACH,
        PASTEL_SAGE,
        PASTEL_PERIWINKLE
    ]

    for idx, series in enumerate(chart.series):
        color = pastel_colors[idx % len(pastel_colors)]
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = color

    return slide


def create_pie_chart_slide(prs, title, categories, values):
    """Create pie chart with pastel colors"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = DARK_SLATE

    # Add chart
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('Values', values)

    x, y, cx, cy = Inches(2), Inches(1.6), Inches(6), Inches(3.8)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    ).chart

    # Style pie chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT

    # Apply pastel colors to slices
    pastel_colors = [
        PASTEL_BLUE,
        PASTEL_MINT,
        PASTEL_LAVENDER,
        PASTEL_PEACH,
        PASTEL_CORAL,
        PASTEL_SAGE
    ]

    for idx, point in enumerate(chart.series[0].points):
        color = pastel_colors[idx % len(pastel_colors)]
        fill = point.format.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Show data labels
    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.number_format = '0%'
    data_labels.font.size = Pt(12)
    data_labels.font.color.rgb = DARK_SLATE

    return slide


def add_speaker_notes(slide, notes_text):
    """Add presenter notes to a slide"""
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text


def main():
    """Main function to create the presentation"""
    print("Creating professional business presentation with pastel colors...")

    # Create presentation (16:9 aspect ratio)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Slide 1: Title Slide
    slide = create_title_slide(
        prs,
        "Q1 2026 Business Review",
        "Quarterly Performance and Strategic Update"
    )
    add_speaker_notes(slide, "Welcome everyone. Today we'll review our Q1 performance and discuss our strategic priorities for Q2.")

    # Slide 2: Agenda
    create_content_slide(prs, "Today's Agenda", [
        (0, "Executive Summary"),
        (0, "Financial Performance"),
        (0, "Key Achievements"),
        (0, "Market Analysis"),
        (0, "Challenges and Solutions"),
        (0, "Q2 Strategic Priorities"),
        (0, "Q&A")
    ])

    # Slide 3: Section Header - Executive Summary
    create_section_header(prs, "Executive Summary")

    # Slide 4: Key Highlights
    create_content_slide(prs, "Q1 2026 Key Highlights", [
        (0, "Revenue Growth: 23% YoY increase"),
        (1, "Exceeded Q1 targets by 8%"),
        (0, "Customer Acquisition: 1,250 new clients"),
        (1, "35% improvement over Q4 2025"),
        (0, "Product Launch: Successfully launched 3 new features"),
        (0, "Team Expansion: Grew team from 145 to 182 employees"),
        (0, "Market Share: Increased from 18% to 22%")
    ])

    # Slide 5: Section Header - Financial Performance
    create_section_header(prs, "Financial Performance")

    # Slide 6: Revenue Trend
    create_bar_chart_slide(
        prs,
        "Revenue Growth Trend ($M)",
        ['Q1 2025', 'Q2 2025', 'Q3 2025', 'Q4 2025', 'Q1 2026'],
        [
            ('Revenue', [12.5, 14.2, 15.8, 18.1, 20.3]),
            ('Target', [12.0, 14.0, 16.0, 18.0, 19.0])
        ]
    )

    # Slide 7: Revenue Distribution
    create_pie_chart_slide(
        prs,
        "Revenue Distribution by Product Line",
        ['Enterprise Solutions', 'SMB Products', 'Cloud Services', 'Professional Services', 'Other'],
        [0.35, 0.25, 0.20, 0.15, 0.05]
    )

    # Slide 8: Financial Metrics Table
    create_table_slide(prs, "Key Financial Metrics", [
        ['Metric', 'Q4 2025', 'Q1 2026', 'Change', 'Target'],
        ['Revenue ($M)', '18.1', '20.3', '+12%', '19.0'],
        ['Gross Margin', '68%', '71%', '+3pts', '70%'],
        ['Operating Income ($M)', '5.4', '6.8', '+26%', '6.0'],
        ['Net Income ($M)', '4.1', '5.2', '+27%', '4.5'],
        ['EBITDA ($M)', '6.2', '7.9', '+27%', '7.0'],
        ['Cash Flow ($M)', '4.8', '6.1', '+27%', '5.5']
    ])

    # Slide 9: Section Header - Key Achievements
    create_section_header(prs, "Key Achievements")

    # Slide 10: Product & Technology
    create_two_column_slide(
        prs,
        "Q1 Major Achievements",
        "Product & Technology",
        [
            "• Launched AI-powered analytics platform",
            "• Released mobile app v3.0 (4.8★ rating)",
            "• Achieved 99.97% uptime",
            "• Reduced API response time by 40%",
            "• Completed security audit (zero critical issues)"
        ],
        "Sales & Marketing",
        [
            "• Expanded into 5 new markets",
            "• Signed 3 enterprise contracts ($5M+ each)",
            "• Increased marketing qualified leads by 65%",
            "• Launched brand refresh campaign",
            "• Won 'Best Innovation' industry award"
        ]
    ])

    # Slide 11: Customer Success
    create_content_slide(prs, "Customer Success Metrics", [
        (0, "Net Promoter Score (NPS): 72 (+8 points)"),
        (1, "Industry leading score"),
        (0, "Customer Retention Rate: 94%"),
        (1, "Up from 91% in Q4 2025"),
        (0, "Average Contract Value: $85K"),
        (1, "18% increase YoY"),
        (0, "Customer Satisfaction (CSAT): 4.6/5.0"),
        (0, "Support Response Time: <2 hours")
    ])

    # Slide 12: Section Header - Challenges
    create_section_header(prs, "Challenges & Solutions")

    # Slide 13: Challenges and Mitigation
    create_two_column_slide(
        prs,
        "Addressing Q1 Challenges",
        "Challenges",
        [
            "• Increased competition in key markets",
            "• Supply chain delays (3-week impact)",
            "• Talent acquisition in specialized roles",
            "• Rising cloud infrastructure costs",
            "• Regulatory changes in EU market"
        ],
        "Solutions Implemented",
        [
            "• Launched competitive differentiation campaign",
            "• Diversified supplier base (4 new partners)",
            "• Enhanced employer brand + referral program",
            "• Optimized cloud architecture (22% cost reduction)",
            "• Established EU compliance task force"
        ]
    ])

    # Slide 14: Section Header - Q2 Priorities
    create_section_header(prs, "Q2 2026 Strategic Priorities")

    # Slide 15: Q2 Priorities
    create_content_slide(prs, "Strategic Focus Areas for Q2", [
        (0, "1. Scale Revenue to $23M (13% growth)"),
        (1, "Focus on enterprise segment expansion"),
        (0, "2. Launch Next-Gen Platform"),
        (1, "Complete beta testing and general release"),
        (0, "3. International Expansion"),
        (1, "Establish presence in APAC region"),
        (0, "4. Team Development"),
        (1, "Hire 25 key positions, launch leadership program"),
        (0, "5. Strategic Partnerships"),
        (1, "Close 2-3 major partnership deals")
    ])

    # Slide 16: Q2 Targets Table
    create_table_slide(prs, "Q2 2026 Targets", [
        ['Goal', 'Target', 'Owner', 'Status'],
        ['Revenue', '$23M', 'Sales', 'On Track'],
        ['New Customers', '1,400', 'Marketing', 'On Track'],
        ['Product Launch', 'May 15', 'Product', 'On Track'],
        ['APAC Expansion', 'June 30', 'Operations', 'Planning'],
        ['Team Hiring', '25 hires', 'HR', 'In Progress'],
        ['Partnership Deals', '2-3 deals', 'BD', 'Negotiating']
    ])

    # Slide 17: Closing Slide
    slide = create_content_slide(prs, "Thank You", [
        (0, "Questions & Discussion"),
        (0, ""),
        (0, "Contact Information:"),
        (1, "Email: [email protected]"),
        (1, "Website: www.company.com"),
        (1, "Follow us: @company")
    ])

    # Save presentation
    output_dir = os.path.join(os.path.dirname(__file__), '../../../../output')
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, 'Professional_Business_Review_Q1_2026.pptx')

    prs.save(output_path)

    print(f"✓ Presentation created successfully!")
    print(f"✓ Location: {output_path}")
    print(f"✓ Total slides: {len(prs.slides)}")
    print(f"✓ Format: PowerPoint (.pptx)")
    print("\nPresentation features:")
    print("  • Modern pastel color scheme")
    print("  • 17 professional slides")
    print("  • Charts: Bar chart, Pie chart")
    print("  • Tables: Financial metrics, Q2 targets")
    print("  • Comparison slides with pastel accents")
    print("  • Speaker notes included")


if __name__ == "__main__":
    main()
