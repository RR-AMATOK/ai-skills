#!/usr/bin/env python3
"""
SAP LeanIX Roles, Responsibilities & Time Investment Presentation Generator
Uses SAP LeanIX template to create a comprehensive professional presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.xmlchemy import OxmlElement
import os

# SAP LeanIX Color Palette (matching their brand)
SAP_BLUE = RGBColor(0, 103, 184)           # Primary SAP blue
SAP_LIGHT_BLUE = RGBColor(92, 186, 255)    # Light accent
SAP_DARK_BLUE = RGBColor(0, 73, 144)       # Dark variant
SAP_TEAL = RGBColor(16, 185, 200)          # Teal accent
SAP_GREEN = RGBColor(48, 180, 97)          # Success/positive
SAP_ORANGE = RGBColor(244, 131, 35)        # Highlight
SAP_PURPLE = RGBColor(145, 67, 255)        # Accent
SAP_GRAY = RGBColor(91, 104, 113)          # Secondary text
DARK_TEXT = RGBColor(50, 50, 50)           # Primary text
LIGHT_GRAY_BG = RGBColor(245, 247, 250)    # Background

def create_title_slide(prs, title, subtitle):
    """Create title slide using template layout"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout

    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = title
        for paragraph in slide.shapes.title.text_frame.paragraphs:
            paragraph.font.size = Pt(44)
            paragraph.font.bold = True
            paragraph.font.color.rgb = SAP_DARK_BLUE

    # Set subtitle
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            shape.text = subtitle
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.size = Pt(24)
                paragraph.font.color.rgb = SAP_GRAY
                paragraph.alignment = PP_ALIGN.CENTER

    return slide

def create_content_slide(prs, title, content_items=None):
    """Create a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Content layout

    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = title
        for paragraph in slide.shapes.title.text_frame.paragraphs:
            paragraph.font.size = Pt(32)
            paragraph.font.bold = True
            paragraph.font.color.rgb = SAP_DARK_BLUE

    # Add content if provided
    if content_items:
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                tf = shape.text_frame
                tf.clear()
                for level, text in content_items:
                    p = tf.paragraphs[0] if len(tf.paragraphs) == 1 and not tf.paragraphs[0].text else tf.add_paragraph()
                    p.text = text
                    p.level = level
                    p.font.size = Pt(18 - (level * 2))
                    p.font.color.rgb = DARK_TEXT
                    p.space_before = Pt(6)
                break

    return slide

def create_table_slide(prs, title, table_data, col_widths=None):
    """Create a slide with a formatted table"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add title
    title_shape = slide.shapes.title
    if not title_shape:
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_shape.text = title
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(32)
        paragraph.font.bold = True
        paragraph.font.color.rgb = SAP_DARK_BLUE

    # Calculate table dimensions
    rows, cols = len(table_data), len(table_data[0])
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(4.0)

    # Add table
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths if provided
    if col_widths:
        for i, width in enumerate(col_widths):
            table.columns[i].width = Inches(width)

    # Fill in table data
    for i, row_data in enumerate(table_data):
        for j, cell_value in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_value)

            # Format header row
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = SAP_BLUE
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
                    paragraph.font.size = Pt(12)
                    paragraph.font.color.rgb = RGBColor(255, 255, 255)
                    paragraph.alignment = PP_ALIGN.LEFT
            else:
                # Alternate row colors for better readability
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_GRAY_BG
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(11)
                    paragraph.font.color.rgb = DARK_TEXT
                    paragraph.alignment = PP_ALIGN.LEFT

            # Set margins
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)

    return slide

def create_bar_chart_slide(prs, title, categories, series_data):
    """Create a slide with a bar chart"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Add title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_shape.text = title
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(32)
        paragraph.font.bold = True
        paragraph.font.color.rgb = SAP_DARK_BLUE

    # Create chart data
    chart_data = CategoryChartData()
    chart_data.categories = categories

    for series_name, values in series_data:
        chart_data.add_series(series_name, values)

    # Add chart
    x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(4)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    return slide

def create_section_header(prs, section_title, subtitle=""):
    """Create a section divider slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[2])  # Section header layout

    if slide.shapes.title:
        slide.shapes.title.text = section_title
        for paragraph in slide.shapes.title.text_frame.paragraphs:
            paragraph.font.size = Pt(44)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.alignment = PP_ALIGN.CENTER

    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = SAP_BLUE

    # Add subtitle if provided
    if subtitle:
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                shape.text = subtitle
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.size = Pt(20)
                    paragraph.font.color.rgb = RGBColor(255, 255, 255)
                    paragraph.alignment = PP_ALIGN.CENTER

    return slide

def main():
    # Load the existing SAP LeanIX presentation as base
    # Note: We'll clear it and rebuild with the new content
    base_path = '../../../template/031922_NYL_Operating_Model.pptx'

    if not os.path.exists(base_path):
        print(f"Error: Base presentation not found at {base_path}")
        # Create new blank presentation
        prs = Presentation()
    else:
        # Load existing presentation to preserve template styling
        prs = Presentation(base_path)

        # Clear existing slides but keep the slide masters
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]

    # Slide 1: Title Slide
    create_title_slide(
        prs,
        "SAP LeanIX Roles, Responsibilities & Time Investment Analysis",
        "Operating Model & Governance Framework"
    )

    # Slide 2: Executive Summary
    content = [
        (0, "9 distinct personas identified across your LeanIX workspace"),
        (0, "Roles span from strategic oversight to tactical data stewardship"),
        (0, "Time investments range from 30 minutes to 12 hours per week"),
        (0, "Clear approval workflows ensure data quality and governance"),
        (0, "Success depends on clear SLAs and bottleneck mitigation")
    ]
    create_content_slide(prs, "Executive Summary", content)

    # Slide 3: Section Header - Role Profiles
    create_section_header(prs, "Role Profiles & Time Investment", "Detailed Analysis of 9 Personas")

    # Slide 4: Enterprise Architect
    table_data = [
        ["Attribute", "Details"],
        ["Fact Sheet Involvement", "Organization (Responsible)"],
        ["Primary Responsibilities", "Governance oversight, cross-domain alignment, strategic architecture decisions"],
        ["Secondary Activities", "Reviewing dashboards, generating reports, facilitating architecture review boards"],
        ["Estimated Time/Week", "4-6 hours"],
        ["Interaction Pattern", "Strategic consumer; heavy report/dashboard usage"]
    ]
    create_table_slide(prs, "1. Enterprise Architect", table_data, [2.5, 6.5])

    # Slide 5: Information Systems Architect
    table_data = [
        ["Attribute", "Details"],
        ["Fact Sheet Involvement", "Application, Systems, Data Object, Interface (all Accountable)"],
        ["Primary Responsibilities", "Final approval authority for application landscape, system integrations, data architecture"],
        ["Secondary Activities", "Reviewing change requests, ensuring alignment, participating in governance"],
        ["Estimated Time/Week", "6-10 hours"],
        ["Interaction Pattern", "Heavy reviewer/approver; key bottleneck risk if understaffed"]
    ]
    create_table_slide(prs, "2. Information Systems Architect", table_data, [2.5, 6.5])

    # Slide 6: Solution Architect
    table_data = [
        ["Attribute", "Details"],
        ["Fact Sheet Involvement", "Application, Systems, Data Object, Interface (all Responsible)"],
        ["Primary Responsibilities", "Day-to-day maintenance of application records, system documentation, interfaces"],
        ["Secondary Activities", "Gathering information from teams, updating lifecycle states, documenting dependencies"],
        ["Estimated Time/Week", "8-12 hours"],
        ["Interaction Pattern", "Primary data steward; highest transactional volume"]
    ]
    create_table_slide(prs, "3. Solution Architect", table_data, [2.5, 6.5])

    # Slide 7: Technology Domain Owner
    table_data = [
        ["Attribute", "Details"],
        ["Fact Sheet Involvement", "IT Component, Tech Category (both Accountable)"],
        ["Primary Responsibilities", "Governance over technology standards, approval of IT components, lifecycle decisions"],
        ["Secondary Activities", "Reviewing Technology Domain SME submissions, aligning with procurement"],
        ["Estimated Time/Week", "3-5 hours"],
        ["Interaction Pattern", "Approval-focused; periodic deep-dives during technology assessments"]
    ]
    create_table_slide(prs, "4. Technology Domain Owner", table_data, [2.5, 6.5])

    # Slide 8: Technology Domain SME
    table_data = [
        ["Attribute", "Details"],
        ["Fact Sheet Involvement", "IT Component (Responsible)"],
        ["Primary Responsibilities", "Maintaining IT component records (versions, vendors, support status)"],
        ["Secondary Activities", "Flagging obsolete technologies, supporting rationalization efforts"],
        ["Estimated Time/Week", "4-6 hours"],
        ["Interaction Pattern", "Technical data steward; feeds information to Domain Owner"]
    ]
    create_table_slide(prs, "5. Technology Domain SME", table_data, [2.5, 6.5])

    # Slide 9: Business Architect
    table_data = [
        ["Attribute", "Details"],
        ["Fact Sheet Involvement", "Business Capability, Business Context, Organization (all Responsible)"],
        ["Primary Responsibilities", "Maintaining business capability maps, business context alignment, org structure"],
        ["Secondary Activities", "Facilitating business-IT alignment workshops, capability assessments"],
        ["Estimated Time/Week", "6-8 hours"],
        ["Interaction Pattern", "Strategic data steward; bridges business and IT perspectives"]
    ]
    create_table_slide(prs, "6. Business Architect", table_data, [2.5, 6.5])

    # Slide 10: Supporting Roles
    table_data = [
        ["Role", "Fact Sheets", "Responsibilities", "Time/Week"],
        ["Business Service Owner", "Business Context", "Ensuring business context accuracy, validating alignment", "1-2 hours"],
        ["Organization Owner", "Organization", "Maintaining org unit records, reporting lines, cost centers", "1-2 hours"],
        ["Observers (3 roles)", "Application", "Consuming information, using reports for decision-making", "0.5-1 hour"]
    ]
    create_table_slide(prs, "7-9. Supporting & Observer Roles", table_data, [2.2, 1.8, 3.5, 1.5])

    # Slide 11: Section Header - Interaction Model
    create_section_header(prs, "Interaction & Approval Flow Model", "Governance & Data Stewardship Patterns")

    # Slide 12: Governance Layers Overview
    content = [
        (0, "Governance & Oversight Layer"),
        (1, "Enterprise Architect - Strategic reporting and cross-domain alignment"),
        (0, "Accountable Layer (Approvers)"),
        (1, "Information Systems Architect - Applications, Systems, Data, Interfaces"),
        (1, "Technology Domain Owner - IT Components, Tech Categories"),
        (0, "Responsible Layer (Data Stewards)"),
        (1, "Solution Architect, Technology Domain SME, Business Architect"),
        (0, "Observer Layer (Consumers)"),
        (1, "Business Owner, Operations Lead, Product/Application Lead")
    ]
    create_content_slide(prs, "Governance Layers Overview", content)

    # Slide 13: Application Change Lifecycle
    content = [
        (0, "Pattern 1: Application Change Lifecycle"),
        (1, "Solution Architect creates/updates Application Fact Sheet"),
        (1, "Information Systems Architect reviews and approves/rejects"),
        (1, "If APPROVED: Observers notified of changes"),
        (1, "If REJECTED: Solution Architect revises and resubmits"),
        (0, ""),
        (0, "Key Success Factor: 48-72 hour approval SLA recommended")
    ]
    create_content_slide(prs, "Approval Workflow - Applications", content)

    # Slide 14: Technology Standards Governance
    content = [
        (0, "Pattern 2: Technology Standards Governance"),
        (1, "Technology Domain SME proposes IT Component (new tool/version)"),
        (1, "Technology Domain Owner reviews and decides"),
        (1, "If APPROVED: Added to Tech Catalog as approved technology"),
        (1, "If REJECTED: Flagged as rationalization candidate for retirement"),
        (0, ""),
        (0, "Benefits: Maintains technology standards, controls sprawl")
    ]
    create_content_slide(prs, "Approval Workflow - Technology", content)

    # Slide 15: Cross-Domain Coordination
    content = [
        (0, "Pattern 3: Cross-Domain Coordination"),
        (1, "Business Architect defines Business Capability"),
        (1, "Links to Applications for business-IT alignment"),
        (1, "Solution Architect validates Application-to-Capability mapping"),
        (1, "Information Systems Architect approves final mapping"),
        (0, ""),
        (0, "Ensures: Business capabilities are accurately linked to supporting applications")
    ]
    create_content_slide(prs, "Approval Workflow - Cross-Domain", content)

    # Slide 16: Section Header - Time Investment
    create_section_header(prs, "Time Investment Analysis", "Annual and Weekly Hour Commitments")

    # Slide 17: Time Investment Summary Table
    table_data = [
        ["Role", "Weekly Hours", "Annual Hours", "Frequency Pattern"],
        ["Enterprise Architect", "4-6 hrs", "200-300 hrs", "Steady with peaks during planning"],
        ["Information Systems Architect", "6-10 hrs", "300-500 hrs", "Consistent; high during transformation"],
        ["Solution Architect", "8-12 hrs", "400-600 hrs", "Highest volume; daily engagement"],
        ["Technology Domain Owner", "3-5 hrs", "150-250 hrs", "Peaks during tech refresh"],
        ["Technology Domain SME", "4-6 hrs", "200-300 hrs", "Steady; increases with new tech"],
        ["Business Architect", "6-8 hrs", "300-400 hrs", "Peaks during strategic planning"],
        ["Business Service Owner", "1-2 hrs", "50-100 hrs", "Periodic validation cycles"],
        ["Organization Owner", "1-2 hrs", "50-100 hrs", "Event-driven (reorgs, M&A)"],
        ["Observers (3 roles)", "0.5-1 hr each", "25-50 hrs each", "Consumption-focused"]
    ]
    create_table_slide(prs, "Time Investment Summary", table_data, [2.5, 1.3, 1.3, 3.9])

    # Slide 18: Time Investment Chart
    categories = ['Ent. Arch', 'IS Arch', 'Sol. Arch', 'Tech Owner', 'Tech SME', 'Bus. Arch', 'BS Owner', 'Org Owner', 'Observers']
    series_data = [
        ('Min Hours/Week', [4, 6, 8, 3, 4, 6, 1, 1, 0.5]),
        ('Max Hours/Week', [6, 10, 12, 5, 6, 8, 2, 2, 1])
    ]
    create_bar_chart_slide(prs, "Weekly Time Investment by Role", categories, series_data)

    # Slide 19: Section Header - Success Factors
    create_section_header(prs, "Key Success Factors", "Critical Elements for Effective Implementation")

    # Slide 20: Success Factors
    content = [
        (0, "Clear Escalation Paths"),
        (1, "Define 48-72 hour SLAs for approval turnaround"),
        (1, "Establish clear escalation procedures for blocked requests"),
        (0, "Bottleneck Mitigation"),
        (1, "Information Systems Architect spans 4 fact sheet types"),
        (1, "Consider deputizing or distributing if volume becomes unmanageable"),
        (0, "Observer Engagement"),
        (1, "Tailor dashboards and reports to Observer needs"),
        (1, "Demonstrate value to drive adoption"),
        (0, "Cadence Alignment"),
        (1, "Synchronize data quality reviews with Architecture Review Boards")
    ]
    create_content_slide(prs, "Key Success Factors", content)

    # Slide 21: Recommendations
    content = [
        (0, "Immediate Actions"),
        (1, "Document approval SLAs and communicate to all stakeholders"),
        (1, "Create role-specific training materials and onboarding guides"),
        (1, "Set up automated notifications for pending approvals"),
        (0, "Short-term (1-3 months)"),
        (1, "Monitor approval bottlenecks and adjust staffing if needed"),
        (1, "Establish regular cadence for data quality reviews"),
        (1, "Build dashboards tailored to each persona's needs"),
        (0, "Long-term (3-6 months)"),
        (1, "Consider LeanIX workflow automations for high-volume processes"),
        (1, "Implement delegation scenarios for coverage during absences")
    ]
    create_content_slide(prs, "Recommendations & Next Steps", content)

    # Slide 22: Q&A
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(1.5))
    title_shape.text = "Questions & Discussion"
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(54)
        paragraph.font.bold = True
        paragraph.font.color.rgb = SAP_BLUE
        paragraph.alignment = PP_ALIGN.CENTER

    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_GRAY_BG

    # Save the presentation
    output_path = '../../../output/SAP_LeanIX_Roles_Analysis.pptx'
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)

    print(f"✓ Presentation created successfully!")
    print(f"✓ Location: {os.path.abspath(output_path)}")
    print(f"✓ Total slides: {len(prs.slides)}")
    print(f"✓ Template: SAP LeanIX")
    print(f"\nPresentation includes:")
    print(f"  • Comprehensive role profiles (9 personas)")
    print(f"  • Detailed responsibility breakdowns")
    print(f"  • Time investment analysis with charts")
    print(f"  • Approval workflow patterns")
    print(f"  • Key success factors and recommendations")

if __name__ == "__main__":
    main()
